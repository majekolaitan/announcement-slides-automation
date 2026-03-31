import os
import re
import json
import zipfile
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

try:
    from pptx.oxml.xmlchemy import OxmlElement
except ImportError:
    from pptx.oxml import OxmlElement

try:
    from pptx.oxml.ns import qn
except ImportError:
    from pptx.oxml import qn

# --- CONSTANTS ---
JSON_FILE = "announcements.json"
TEMPLATE_FILE = "SDA_Template.pptx"
OUTPUT_FILE = "SDA_Kubwa_Announcements.pptx"

def clean_json_string(raw_text):
    """Extracts only the JSON array and fixes trailing commas."""
    # Find everything from the first '[' to the last ']'
    match = re.search(r'\[.*\]', raw_text, re.DOTALL)
    if not match:
        raise ValueError("No JSON array (started with '[') found in the text.")
    
    cleaned = match.group(0)
    # Fix trailing commas
    cleaned = re.sub(r',\s*([\]}])', r'\1', cleaned)
    return cleaned

def inject_loop_into_zip(filename):
    """Hacks the PPTX XML to make the presentation loop automatically."""
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(filename, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        presProps_path = os.path.join(temp_dir, 'ppt', 'presProps.xml')
        if os.path.exists(presProps_path):
            with open(presProps_path, 'r', encoding='utf-8') as f:
                presProps = f.read()
            
            # Inject loop and showNarration parameters
            presProps = re.sub(
                r"(<p:presentationPr[^>]*>)", 
                r'\1<p:showPr showNarration="1" loop="1"/>', 
                presProps, 
                count=1
            )
            
            with open(presProps_path, 'w', encoding='utf-8') as f:
                f.write(presProps)
        
        # Repackage the zip
        with zipfile.ZipFile(filename, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zip_ref.write(file_path, arcname)

def main():
    # 1. Load and parse JSON
    if not os.path.exists(JSON_FILE):
        print(f"Error: '{JSON_FILE}' missing. Please create it and paste the AI output inside.")
        return

    with open(JSON_FILE, 'r', encoding='utf-8') as f:
        raw_data = f.read()

    try:
        clean_data = clean_json_string(raw_data)
        events_data = json.loads(clean_data)
    except json.JSONDecodeError as e:
        print(f"Error parsing JSON. Please ensure the AI output is valid. Details: {e}")
        return

    # 2. Load Template OR Fallback to Default
    if os.path.exists(TEMPLATE_FILE):
        print(f"Applying custom template: {TEMPLATE_FILE}")
        prs = Presentation(TEMPLATE_FILE)
    else:
        print(f"Template '{TEMPLATE_FILE}' not found. Falling back to default 16:9 layout.")
        prs = Presentation()
        # Set to standard widescreen (16:9) aspect ratio
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    # 3. Generate slides based on data
    for event in events_data:
        # Fallback to layout 0 if layout 1 is missing
        try:
            slide_layout = prs.slide_layouts[1]
        except IndexError:
            slide_layout = prs.slide_layouts[0]
            
        slide = prs.slides.add_slide(slide_layout)

        # --- Title Configuration ---
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text_frame.clear()
            p_title = title_shape.text_frame.paragraphs[0]
            
            # Note: We enforce LEFT alignment and BOLD here. 
            # Font size/family will inherit from the Slide Master!
            p_title.alignment = PP_ALIGN.LEFT
            run_title = p_title.add_run()
            run_title.text = event.get("title", "Announcement")
            run_title.font.bold = True

        # --- Body (Bullets) Configuration ---
        try:
            # placeholders[1] is the standard body text box in Layout 1
            body_shape = slide.placeholders[1]
            body_shape.text_frame.clear()
            
            for i, bullet in enumerate(event.get("bullets", [])):
                if i == 0:
                    p_body = body_shape.text_frame.paragraphs[0]
                else:
                    p_body = body_shape.text_frame.add_paragraph()
                    
                run_body = p_body.add_run()
                run_body.text = bullet
                run_body.font.bold = True
                # No font size or name is defined here so it perfectly 
                # matches whatever is set in your Slide Master.
        except KeyError:
            print("Warning: Slide Master Layout 1 is missing a body placeholder.")

        # --- Main Icon Configuration ---
        # Dynamic Main Icon positioning (Adapts to 16:9 or 4:3 templates!)
        icon_width = Inches(2.0)
        icon_height = Inches(1.5)
        icon_left = prs.slide_width - icon_width - Inches(0.5) # 0.5 inches from the right edge
        icon_top = Inches(0.5) # 0.5 inches from the top
        
        icon_box = slide.shapes.add_textbox(icon_left, icon_top, icon_width, icon_height)
        icon_frame = icon_box.text_frame
        icon_frame.clear()
        p_icon = icon_frame.paragraphs[0]
        p_icon.alignment = PP_ALIGN.RIGHT
        run_icon = p_icon.add_run()
        run_icon.text = event.get("icon", "📌")
        run_icon.font.size = Pt(72)
        run_icon.font.name = 'Segoe UI Emoji'

        # --- Transition Configuration (Push Left, Auto-Advance 8s) ---
        transition = slide.element.find(qn('p:transition'))
        if transition is not None:
            transition.getparent().remove(transition)
            
        new_transition = OxmlElement('p:transition')
        new_transition.set('advTm', '8000') # 8 seconds
        push = OxmlElement('p:push')
        push.set('dir', 'l')
        new_transition.append(push)
        
        inserted = False
        for tag in ['p:timing', 'p:extLst']:
            elem = slide.element.find(qn(tag))
            if elem is not None:
                elem.addprevious(new_transition)
                inserted = True
                break
        if not inserted:
            slide.element.append(new_transition)

    # 4. Save and Apply Loop Hack
    try:
        prs.save(OUTPUT_FILE)
        try:
            inject_loop_into_zip(OUTPUT_FILE)
            print(f"Success! '{OUTPUT_FILE}' generated and set to loop automatically.")
        except Exception as e:
            print(f"Slides generated, but looping could not be applied: {e}")
    except PermissionError:
        print(f"Error: Permission denied. Close '{OUTPUT_FILE}' if it is open and try again.")

if __name__ == "__main__":
    main()